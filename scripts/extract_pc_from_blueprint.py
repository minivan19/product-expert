#!/usr/bin/env python3
"""
product-expert 场景四 - 从蓝图方案提取产品方案卡

用法:
  python3 scripts/extract_pc_from_blueprint.py --customer "诺斯贝尔" --blueprint "路径"
  python3 scripts/extract_pc_from_blueprint.py --customer "诺斯贝尔" --auto  # 自动找最新蓝图

场景四a: 蓝图方案 → LLM凝练 → 提取产品方案卡（一对多）
  输入: 客户蓝图方案文件(docx/pdf/pptx)
  处理: LLM分析蓝图中的产品实施内容 → 凝练为产品方案卡
  输出: 产品方案卡JSON草案，供CSM确认

场景四b: 场景三方案 → 询问是否入库
  输入: 场景三生成的实施路线图JSON
  处理: 解析处理方式 → 映射产品模块 → 生成方案卡草案
  输出: 询问CSM是否保存为正式卡片

场景四c: 卡片库Review
  定期扫描，识别空白产品模块，提示补充
"""

import json
import sys
import os
import subprocess
from pathlib import Path
from datetime import datetime

# ── 路径配置 ──────────────────────────────────────────────

SKILL_DIR = Path(__file__).parent.parent
PRODUCT_CARDS_DIR = SKILL_DIR / "framework" / "product_solution_cards"
CARD_INDEX_FILE = SKILL_DIR / "framework" / "product_card_index.json"
PC_SCHEMA_FILE = SKILL_DIR / "framework" / "product_solution_card_schema.json"
CLIENT_DATA_ROOT = Path("/Users/limingheng/AI/client-data")

# ── 工具函数 ──────────────────────────────────────────────

def load_json(path):
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)

def save_json(path, data):
    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)

def get_next_card_id() -> str:
    """获取下一个PC_XX ID"""
    existing = list(PRODUCT_CARDS_DIR.glob("PC_*.json"))
    if not existing:
        return "PC_01"
    nums = []
    for f in existing:
        try:
            num = int(f.stem.split("_")[1])
            nums.append(num)
        except:
            pass
    return f"PC_{max(nums) + 1:02d}"

def read_blueprint(file_path: str) -> str:
    """读取蓝图文件内容，支持docx/pdf/pptx"""
    ext = Path(file_path).suffix.lower()
    text = ""

    if ext == ".docx":
        try:
            from docx import Document
            doc = Document(file_path)
            text = "\n".join([p.text for p in doc.paragraphs])
            # 读表格
            for tbl in doc.tables:
                for row in tbl.rows:
                    for cell in row.cells:
                        text += "\n" + cell.text
        except ImportError:
            text = f"[需要python-docx库读取docx文件: {file_path}]"
    elif ext == ".pdf":
        text = f"[PDF文件需要单独处理: {file_path}]"
    elif ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += "\n" + shape.text
        except ImportError:
            text = f"[需要python-pptx库读取pptx文件: {file_path}]"
        text = f"[PPTX文件需要单独处理: {file_path}]"
    else:
        # 尝试直接读
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
        except:
            text = f"[无法读取文件: {file_path}]"

    return text.strip()

def find_latest_blueprint(customer_name: str) -> Path:
    """自动找客户最新蓝图方案"""
    customer_dir = CLIENT_DATA_ROOT / "客户档案" / customer_name / "蓝图方案"
    if not customer_dir.exists():
        # 尝试上级目录
        customer_dir = CLIENT_DATA_ROOT / "客户档案" / customer_name

    candidates = list(customer_dir.glob("**/*蓝图*.docx")) + \
                 list(customer_dir.glob("**/*蓝图*.pptx")) + \
                 list(customer_dir.glob("**/*蓝图*.pdf"))

    if candidates:
        # 返回最新的
        return max(candidates, key=lambda p: p.stat().st_mtime)
    return None

# ── LLM 提取 ──────────────────────────────────────────────

DEEPSEEK_API_KEY = os.environ.get("DEEPSEEK_API_KEY", "")
DEEPSEEK_API_KEY = DEEPSEEK_API_KEY or os.environ.get("OPENAI_API_KEY", "")

def llm_extract_cards(blueprint_text: str, customer_name: str) -> list:
    """
    用DeepSeek LLM从蓝图文本中提取产品方案卡。
    返回: [{"name": "...", "modules": [...], "steps": [...], "suitable_for": {...}}, ...]
    """
    if not blueprint_text or len(blueprint_text) < 100:
        return [{"error": "蓝图内容太短或无法读取"}]

    # 控制token：只取前8000字
    text_to_send = blueprint_text[:8000]

    prompt = f"""你是一个ERP实施专家。请从以下客户蓝图方案中提取「产品方案卡」。

产品方案卡格式说明：
- name: 方案名称（简洁，如"供应商准入+绩效管理联动方案"）
- modules: 产品模块组合列表，如["基础供应商管理", "高级供应商管理"]
- business_domains: 对应的业务域ID列表，如["DM_07_01", "DM_07_03"]
- suitable_for: 适用条件 {{行业: [...], 物料属性: [...], 权力结构: [...], 协同成熟度: [...]}}
- steps: 实施步骤列表，每步为字符串
- config_items: 关键配置项列表
- customization_points: 二开介入点描述列表（没有则填"无"）

注意：
1. 一份蓝图可能包含多个产品方案卡，请尽可能拆分
2. 只提取蓝图中有明确产品对应的内容，不要推理
3. 如果某节内容不涉及具体产品功能，跳过
4. 二开介入点指现有产品无法直接满足、需要定制开发的功能点

客户名称：{customer_name}
蓝图内容：
---
{text_to_send}
---

请以JSON数组格式输出产品方案卡列表，key用英文。
"""

    import urllib.request
    import urllib.error

    payload = {
        "model": "deepseek-chat",
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0.3,
        "max_tokens": 2000
    }

    req = urllib.request.Request(
        "https://api.deepseek.com/chat/completions",
        data=json.dumps(payload).encode("utf-8"),
        headers={
            "Authorization": f"Bearer {DEEPSEEK_API_KEY}",
            "Content-Type": "application/json"
        },
        method="POST"
    )

    try:
        with urllib.request.urlopen(req, timeout=60) as resp:
            result = json.loads(resp.read().decode("utf-8"))
            content = result["choices"][0]["message"]["content"]
            # 尝试解析JSON
            if "```json" in content:
                start = content.find("```json") + 7
                end = content.find("```", start)
                content = content[start:end]
            elif "```" in content:
                start = content.find("```") + 3
                end = content.find("```", start)
                content = content[start:end]
            return json.loads(content.strip())
    except Exception as e:
        return [{"error": str(e)}]


def build_pc_from_llm_result(llm_card: dict, customer_name: str, source_file: str) -> dict:
    """将LLM提取结果构建为标准产品方案卡"""
    card_id = get_next_card_id()
    now = datetime.now().isoformat()

    card = {
        "id": card_id,
        "名称": llm_card.get("name", "未命名方案"),
        "版本": "v1",
        "状态": "待确认",
        "来源": {
            "类型": "蓝图提取",
            "来源客户": customer_name,
            "来源文件": source_file,
            "提取时间": now
        },
        "对应业务卡": [],  # 待关联
        "适用条件": {
            "行业": llm_card.get("suitable_for", {}).get("行业", []),
            "物料属性": llm_card.get("suitable_for", {}).get("物料属性", []),
            "权力结构": llm_card.get("suitable_for", {}).get("权力结构", []),
            "协同成熟度": llm_card.get("suitable_for", {}).get("协同成熟度", [])
        },
        "产品模块组合": [
            {
                "模块": m,
                "套件": "",  # 待填充
                "角色": "主要",
                "配置项": [],
                "二开介入": False
            }
            for m in llm_card.get("modules", [])
        ],
        "实施步骤": [
            {"顺序": i+1, "步骤": s, "产品模块": "", "前置依赖": "", "二开": False}
            for i, s in enumerate(llm_card.get("steps", []))
        ],
        "二开介入点": [
            {"环节": pt, "描述": "", "难度": "中", "预估人天": ""}
            for pt in llm_card.get("customization_points", []) if pt != "无"
        ],
        "关键配置项": llm_card.get("config_items", []),
        "成功客户": [{"客户": customer_name, "实施年份": datetime.now().year, "效果": ""}],
        "关联方案卡": [],
        "更新时间": now,
        "维护人": ""
    }
    return card


# ── 场景四a: 蓝图提取 ─────────────────────────────────────

def scene_4a_blueprint(customer_name: str, blueprint_path: str = None, auto: bool = False):
    """从蓝图方案提取产品方案卡"""
    print(f"\n🌐 场景四a: 蓝图方案 → 产品方案卡")
    print(f"客户: {customer_name}")

    # 找蓝图文件
    if auto:
        bp = find_latest_blueprint(customer_name)
        if bp:
            blueprint_path = str(bp)
            print(f"  → 自动找到: {bp.name}")
        else:
            print(f"  ⚠️ 未找到蓝图文件")
            return
    elif blueprint_path:
        print(f"  → 指定文件: {blueprint_path}")
    else:
        print(f"  ⚠️ 请提供蓝图文件路径或使用 --auto")
        return

    # 读取内容
    text = read_blueprint(blueprint_path)
    if "error" in text.lower() or len(text) < 100:
        print(f"  ⚠️ 文件无法读取或内容太短: {text[:200]}")
        return

    print(f"  → 蓝图内容长度: {len(text)} 字")

    # LLM提取
    print(f"  → 调用LLM提取产品方案卡...")
    cards = llm_extract_cards(text, customer_name)

    if not cards or (len(cards) == 1 and "error" in cards[0]):
        print(f"  ⚠️ LLM提取失败: {cards[0].get('error', '未知错误')}")
        return

    print(f"  → 提取到 {len(cards)} 张产品方案卡:")

    saved = []
    for i, llm_card in enumerate(cards, 1):
        card = build_pc_from_llm_result(llm_card, customer_name, blueprint_path)
        print(f"\n  【{card['id']}】{card['名称']}")
        print(f"    产品模块: {[m['模块'] for m in card['产品模块组合']]}")
        print(f"    实施步骤: {len(card['实施步骤'])} 步")

        # 询问是否保存
        save = input(f"    保存此卡片? [Y/n]: ").strip().lower()
        if save != "n":
            card_file = PRODUCT_CARDS_DIR / f"{card['id']}_{llm_card.get('name', 'card')[:20]}.json"
            save_json(card_file, card)
            saved.append(card)
            print(f"    ✅ 已保存: {card_file.name}")

    print(f"\n✅ 共保存 {len(saved)} 张产品方案卡")

    # 更新索引
    update_card_index(saved)

    return saved


# ── 场景四b: 场景三方案入库 ───────────────────────────────

def scene_4b_save_from_scene3(scene3_json_path: str):
    """
    从场景三输出的实施路线图JSON中提取方案卡草案，
    询问CSM是否保存。
    """
    print(f"\n🌐 场景四b: 场景三方案 → 产品方案卡")

    try:
        data = json.loads(Path(scene3_json_path).read_text())
    except:
        print(f"  ⚠️ 无法读取JSON文件: {scene3_json_path}")
        return

    # 解析场景三输出
    solution = data.get("方案", {})
    card_name = solution.get("方案名称", "未命名")
    methods = [m["处理方式"] for m in solution.get("推荐处理方式", [])]

    print(f"  方案名称: {card_name}")
    print(f"  处理方式: {methods}")

    confirm = input(f"\n  保存为产品方案卡? [y/N]: ").strip().lower()
    if confirm != "y":
        print("  → 取消")
        return

    card_id = get_next_card_id()
    now = datetime.now().isoformat()

    card = {
        "id": card_id,
        "名称": card_name,
        "版本": "v1",
        "状态": "草稿",
        "来源": {
            "类型": "场景三生成",
            "来源客户": solution.get("客户", ""),
            "提取时间": now
        },
        "对应业务卡": [],  # 待关联
        "适用条件": {
            "行业": [],
            "物料属性": [],
            "权力结构": [],
            "协同成熟度": []
        },
        "产品模块组合": [],  # 待填充（需要Qdrant或手动）
        "实施步骤": [],
        "二开介入点": [],
        "关键配置项": [],
        "成功客户": [],
        "关联方案卡": [],
        "更新时间": now,
        "维护人": ""
    }

    print(f"\n  【{card_id}】{card['名称']} (草稿状态)")
    print(f"  → 需要后续补充: 产品模块组合、实施步骤、适用条件")
    print(f"  → 状态为'草稿'，需CSM确认后转为正式卡片")

    card_file = PRODUCT_CARDS_DIR / f"{card_id}_{card_name[:20]}.json"
    save_json(card_file, card)
    print(f"  ✅ 已保存: {card_file.name}")
    update_card_index([card])


# ── 场景四c: 卡片库Review ────────────────────────────────

def scene_4c_review():
    """卡片库定期Review"""
    print(f"\n🌐 场景四c: 产品方案卡库Review")

    index = load_json(CARD_INDEX_FILE)
    print(f"  总卡数: {index['汇总']['总卡数']}")
    print(f"  来源分布: {index['汇总']['来源分布']}")

    # 统计各业务域覆盖
    domain_coverage = {}
    for entry in index.get("索引", []):
        domain = entry.get("域", "未知")
        cards = entry.get("已有卡片", [])
        domain_coverage[domain] = cards

    print(f"\n  业务域覆盖:")
    for domain, cards in domain_coverage.items():
        status = "✅" if cards else "⚠️ 空白"
        print(f"    {status} {domain}: {cards or '无卡片'}")

    # 提示空白模块
    blank_modules = index.get("汇总", {}).get("空白产品模块", [])
    if blank_modules:
        print(f"\n  建议补充（空白产品模块）:")
        for m in blank_modules[:5]:
            print(f"    • {m}")

    return index


# ── 索引更新 ──────────────────────────────────────────────

def update_card_index(new_cards: list):
    """更新产品方案卡索引"""
    index = load_json(CARD_INDEX_FILE)

    index["汇总"]["总卡数"] = len(list(PRODUCT_CARDS_DIR.glob("PC_*.json")))

    for card in new_cards:
        src_type = card.get("来源", {}).get("类型", "手工创建")
        if src_type == "蓝图提取":
            index["汇总"]["来源分布"]["蓝图提取"] += 1
        elif src_type == "场景三生成":
            index["汇总"]["来源分布"]["场景三生成"] += 1
        else:
            index["汇总"]["来源分布"]["手工创建"] += 1

    save_json(CARD_INDEX_FILE, index)


# ── CLI ──────────────────────────────────────────────────

def main():
    import argparse
    parser = argparse.ArgumentParser(description="Product Expert 场景四 - 产品方案卡管理")
    parser.add_argument("--customer", "-c", help="客户名称（场景四a）")
    parser.add_argument("--blueprint", "-b", help="蓝图文件路径（场景四a）")
    parser.add_argument("--auto", "-a", action="store_true", help="自动找最新蓝图（场景四a）")
    parser.add_argument("--scene3-json", help="场景三JSON路径（场景四b）")
    parser.add_argument("--review", "-r", action="store_true", help="场景四c: 卡片库Review")
    parser.add_argument("--list", "-l", action="store_true", help="列出所有产品方案卡")

    args = parser.parse_args()

    if args.list:
        cards = list(PRODUCT_CARDS_DIR.glob("PC_*.json"))
        print(f"\n📦 产品方案卡库 ({len(cards)} 张):")
        for c in sorted(cards):
            d = load_json(c)
            print(f"  [{d.get('id', c.stem)}] {d.get('名称', '?')}")
            print(f"       状态:{d.get('状态','?')} | 来源:{d.get('来源',{}).get('类型','?')} | 客户:{d.get('来源',{}).get('来源客户','?')}")
        return

    if args.review:
        scene_4c_review()
        return

    if args.scene3_json:
        scene_4b_save_from_scene3(args.scene3_json)
        return

    if args.customer:
        scene_4a_blueprint(args.customer, args.blueprint, args.auto)
        return

    parser.print_help()


if __name__ == "__main__":
    main()
